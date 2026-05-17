"""Construit le PPTX de soutenance de Thèse Professionnelle (~24 slides).

Présentation EXHAUSTIVE : couvre tout le plan du mémoire déposé, section par
section (Introduction · Ch.1 Revue de littérature 1.1-1.5 · Ch.2 Concepts/méthodes
2.1-2.6 · Ch.3 Résultats 3.1-3.3 · Conclusion).

Style institutionnel sobre : titres serif, corps sans-serif, bandeau de section,
pied de page (nom + n° de slide). Le contenu reprend fidèlement le mémoire ; les
chiffres proviennent de `evals/summary.json`.

Mise en forme et orchestration uniquement — aucune logique du pipeline n'est
réimplémentée.

Pré-requis : lancer d'abord `python presentation/generate_figures.py`.

Usage:
    python presentation/build_deck.py
"""

from __future__ import annotations

import json
from pathlib import Path

import numpy as np
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt

ROOT = Path(__file__).resolve().parent.parent
HERE = Path(__file__).resolve().parent
FIGURES = HERE / "figures"
SUMMARY = ROOT / "evals" / "summary.json"
OUT = HERE / "soutenance.pptx"

# Palette institutionnelle sobre
NAVY = RGBColor(0x1B, 0x2A, 0x4A)
GREY = RGBColor(0x5D, 0x6D, 0x7E)
LIGHT = RGBColor(0xB5, 0xBE, 0xC9)
ACCENT = RGBColor(0x2E, 0x86, 0xC1)
DARK = RGBColor(0x2B, 0x2B, 0x2B)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
BLUEGREY = RGBColor(0x5D, 0x72, 0x99)

SERIF = "Georgia"
SANS = "Calibri"

EMU_W, EMU_H = Inches(13.333), Inches(7.5)
FOOTER_TEXT = "M. Ennagi — Soutenance de Thèse Professionnelle · Juin 2026"
METRICS = ["faithfulness", "answer_relevancy", "context_precision", "context_recall"]


# --------------------------------------------------------------------------- #
# Données réelles
# --------------------------------------------------------------------------- #
def load_canonical_runs() -> list[dict]:
    runs = []
    with SUMMARY.open() as f:
        for line in f:
            line = line.strip()
            if line and "embedding_model" in line:
                runs.append(json.loads(line))
    return runs


def config_label(row: dict) -> str:
    hyde, rerank = row.get("hyde", False), row.get("rerank", False)
    if hyde and rerank:
        return "hyde+rerank"
    if rerank:
        return "+rerank"
    if hyde:
        return "+hyde"
    return "baseline"


def runs_for_model(runs: list[dict], model: str) -> dict[str, dict]:
    return {config_label(r): r["scores"] for r in runs
            if r.get("generation_model") == model}


def scores_table_12(runs: list[dict]) -> list[list[str]]:
    """Score moyen (moyenne des 4 métriques) par modèle × configuration."""
    configs = ["baseline", "+rerank", "+hyde", "hyde+rerank"]
    models = [m for m in ["qwen2.5:7b", "mistral:7b", "gemma2:9b"] if runs_for_model(runs, m)]
    rows = [["Modèle"] + configs]
    for model in models:
        data = runs_for_model(runs, model)
        line = [model]
        for c in configs:
            line.append(f"{np.mean([data[c][m] for m in METRICS]):.3f}" if c in data else "—")
        rows.append(line)
    return rows


# --------------------------------------------------------------------------- #
# Helpers de mise en page
# --------------------------------------------------------------------------- #
def blank_slide(prs: Presentation):
    return prs.slides.add_slide(prs.slide_layouts[6])


def set_bg(slide, color=WHITE):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_textbox(slide, left, top, width, height, anchor=MSO_ANCHOR.TOP):
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    return tf


def style_run(run, size, color=DARK, bold=False, italic=False, font=SANS):
    run.font.size = Pt(size)
    run.font.color.rgb = color
    run.font.bold = bold
    run.font.italic = italic
    run.font.name = font


def add_line(tf, text, size, color=DARK, bold=False, italic=False, font=SANS,
             align=None, first=False, space_after=None, space_before=None):
    p = tf.paragraphs[0] if first else tf.add_paragraph()
    if align is not None:
        p.alignment = align
    if space_after is not None:
        p.space_after = Pt(space_after)
    if space_before is not None:
        p.space_before = Pt(space_before)
    r = p.add_run()
    r.text = text
    style_run(r, size, color, bold=bold, italic=italic, font=font)
    return p


def hline(slide, left, top, width, color=NAVY, weight=1.25):
    ln = slide.shapes.add_connector(2, left, top, left + width, top)
    ln.line.color.rgb = color
    ln.line.width = Pt(weight)
    return ln


def header(slide, section, title):
    sec = add_textbox(slide, Inches(0.7), Inches(0.32), Inches(12), Inches(0.32))
    add_line(sec, section.upper(), 12, GREY, bold=True, font=SANS, first=True)
    ttl = add_textbox(slide, Inches(0.7), Inches(0.64), Inches(12), Inches(0.7))
    add_line(ttl, title, 24, NAVY, bold=True, font=SERIF, first=True)
    hline(slide, Inches(0.7), Inches(1.38), Inches(11.93), color=NAVY, weight=1.25)


def footer(slide, n, total):
    hline(slide, Inches(0.7), Inches(7.05), Inches(11.93), color=LIGHT, weight=0.75)
    left = add_textbox(slide, Inches(0.7), Inches(7.08), Inches(9), Inches(0.35))
    add_line(left, FOOTER_TEXT, 9, GREY, font=SANS, first=True)
    right = add_textbox(slide, Inches(11.5), Inches(7.08), Inches(1.13), Inches(0.35))
    add_line(right, f"{n} / {total}", 9, GREY, font=SANS, align=PP_ALIGN.RIGHT, first=True)


def bullets(slide, items, left=Inches(0.7), top=Inches(1.65),
            width=Inches(12), height=Inches(5.1), size=16):
    tf = add_textbox(slide, left, top, width, height)
    for i, item in enumerate(items):
        text, level = (item if isinstance(item, tuple) else (item, 0))
        bullet = "•  " if level == 0 else "–  "
        color = DARK if level == 0 else GREY
        p = add_line(tf, bullet + text, size - level, color, font=SANS,
                     first=(i == 0), space_after=6)
        p.level = level
    return tf


def notes(slide, text):
    slide.notes_slide.notes_text_frame.text = " ".join(text.split())


def card(slide, left, top, width, height, titre, txt, color):
    b = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    b.fill.solid()
    b.fill.fore_color.rgb = color
    b.line.fill.background()
    tf = b.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.TOP
    tf.margin_left = Inches(0.2)
    tf.margin_right = Inches(0.2)
    tf.margin_top = Inches(0.18)
    add_line(tf, titre, 15, WHITE, bold=True, font=SERIF, align=PP_ALIGN.CENTER, first=True)
    add_line(tf, txt, 12, WHITE, font=SANS, align=PP_ALIGN.CENTER, space_before=8)
    return b


def arrow(slide, left, top, width=Inches(0.3), height=Inches(0.45)):
    a = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, left, top, width, height)
    a.fill.solid()
    a.fill.fore_color.rgb = LIGHT
    a.line.fill.background()
    return a


def code_box(slide, left, top, width, height, code, size=11, title=None):
    """Affiche du code (monospace) sur fond gris clair."""
    box = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    box.fill.solid()
    box.fill.fore_color.rgb = RGBColor(0xF2, 0xF4, 0xF7)
    box.line.color.rgb = LIGHT
    box.line.width = Pt(0.75)
    tf = box.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.TOP
    tf.margin_left = Inches(0.18)
    tf.margin_right = Inches(0.1)
    tf.margin_top = Inches(0.12)
    first = True
    if title:
        add_line(tf, title, size, GREY, italic=True, font="Consolas", first=True)
        first = False
    for ln in code.split("\n"):
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        if "#" in ln:  # séparer code (navy) et commentaire (gris)
            head, _, comment = ln.partition("#")
            r = p.add_run()
            r.text = head
            style_run(r, size, NAVY, font="Consolas")
            r2 = p.add_run()
            r2.text = "#" + comment
            style_run(r2, size, GREY, italic=True, font="Consolas")
        else:
            r = p.add_run()
            r.text = ln if ln else " "
            style_run(r, size, NAVY, font="Consolas")
    return box


# --------------------------------------------------------------------------- #
# Slides — partie liminaire
# --------------------------------------------------------------------------- #
def slide_titre(prs):
    s = blank_slide(prs)
    hline(s, Inches(1), Inches(1.2), Inches(11.33), color=NAVY, weight=2)
    tf = add_textbox(s, Inches(1), Inches(1.7), Inches(11.33), Inches(2.6), MSO_ANCHOR.MIDDLE)
    add_line(tf, "Conception et évaluation d'une architecture RAG souveraine pour la "
             "capitalisation des connaissances et l'optimisation de l'accès documentaire",
             26, NAVY, bold=True, font=SERIF, align=PP_ALIGN.CENTER, first=True)
    add_line(tf, "Cas du Fonds de l'Accompagnement des Réformes du Transport Routier (FART)",
             16, GREY, italic=True, font=SERIF, align=PP_ALIGN.CENTER, space_before=10)
    hline(s, Inches(1), Inches(4.5), Inches(11.33), color=NAVY, weight=2)
    meta = add_textbox(s, Inches(1), Inches(4.8), Inches(11.33), Inches(2.0), MSO_ANCHOR.MIDDLE)
    rows = [("Thèse Professionnelle — [Programme / Master]", 15, True),
            ("Présentée par : [Votre Nom]", 14, False),
            ("Directeur de thèse : [Nom] · Jury : [Membres]", 13, False),
            ("[Établissement] — Juin 2026", 13, False)]
    for i, (txt, sz, bold) in enumerate(rows):
        add_line(meta, txt, sz, NAVY if bold else GREY, bold=bold, font=SANS,
                 align=PP_ALIGN.CENTER, first=(i == 0), space_after=4)
    notes(s, """Madame, Monsieur les membres du jury, je vous présente ma thèse
    professionnelle : la conception et l'évaluation d'une architecture RAG souveraine pour la
    capitalisation des connaissances, appliquée au cas du FART. ÉDITEZ les placeholders.""")


def slide_plan(prs):
    s = blank_slide(prs)
    header(s, "Soutenance", "Plan de la présentation")
    bullets(s, [
        "Introduction — contexte, problème et problématique du FART",
        "Questions de recherche et hypothèses",
        "Chapitre 1 — Revue de la littérature (1.1 à 1.5)",
        "Chapitre 2 — Concepts, méthodes et outils mobilisés (2.1 à 2.6)",
        "Chapitre 3 — Présentation et discussion des résultats (3.1 à 3.3)",
        "Discussion, transposition au FART et conclusion générale",
    ], size=18, top=Inches(1.9))
    notes(s, """Le plan de ma présentation suit fidèlement celui du mémoire : introduction et
    problématique, puis les trois chapitres — revue de la littérature, concepts et méthodes,
    résultats — et enfin la discussion et la conclusion.""")


def slide_intro(prs):
    s = blank_slide(prs)
    header(s, "Introduction", "Contexte : la mémoire institutionnelle du FART")
    bullets(s, [
        "La transformation numérique a sécurisé la conservation des données mais engendré une "
        "surcharge informationnelle.",
        "Le FART administre un patrimoine documentaire hétérogène et critique — sa mémoire "
        "institutionnelle :",
        ("normes juridiques (lois de finances, décrets, arrêtés) ;", 1),
        ("documents procéduraux (manuels, guides d'éligibilité) ;", 1),
        ("documents décisionnels (procès-verbaux, notes de synthèse).", 1),
        "Patrimoine fragmenté : arborescences cloisonnées, PDF non structurés, scans.",
    ])
    notes(s, """Le contexte. La dématérialisation a produit une surcharge informationnelle. Le
    FART, acteur de la mobilité urbaine, détient un patrimoine hétérogène et critique — sa
    mémoire institutionnelle — composé de normes juridiques, de documents procéduraux et
    décisionnels, mais fragmenté et souvent sous forme de scans non structurés.""")


def slide_probleme(prs):
    s = blank_slide(prs)
    header(s, "Introduction", "Problème central et contraintes")
    bullets(s, [
        "Aucun moteur d'indexation centralisé : exploration manuelle d'arborescences.",
        "Archivage individuel en silos → patrimoine invisible ; recherche sémantique impossible.",
        "Dépendance au savoir tacite : la mobilité du personnel érode la mémoire collective.",
        "Conséquences : latence de traitement, risque de discordance réglementaire.",
    ], top=Inches(1.6), height=Inches(2.4))
    top = Inches(4.3)
    w = Inches(3.9)
    gap = Inches(0.2)
    cards = [
        ("Souveraineté", "Pas de SaaS : architecture strictement locale (On-Premise).", NAVY),
        ("Risque hallucinatoire", "Fidélité documentaire stricte, restitution sourcée.", ACCENT),
        ("Hétérogénéité", "Documents numérisés : chaîne OCR et segmentation.", GREY),
    ]
    for i, (t, txt, c) in enumerate(cards):
        card(s, Inches(0.7) + i * (w + gap), top, w, Inches(2.1), t, txt, c)
    notes(s, """Le diagnostic : pas d'indexation centralisée, des silos individuels, une
    recherche sémantique impossible, une dépendance au savoir tacite. Résultat : lenteur et
    risque de discordance. Trois contraintes bloquantes encadrent la solution : la souveraineté
    (tout en local), le risque hallucinatoire (fidélité stricte, sources citées), et
    l'hétérogénéité des documents numérisés (OCR et segmentation).""")


def slide_problematique(prs):
    s = blank_slide(prs)
    header(s, "Problématique", "Question de recherche principale")
    quote = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                               Inches(1.0), Inches(2.2), Inches(11.33), Inches(2.6))
    quote.fill.solid()
    quote.fill.fore_color.rgb = RGBColor(0xF2, 0xF4, 0xF7)
    quote.line.color.rgb = NAVY
    quote.line.width = Pt(1)
    tf = quote.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    tf.margin_left = Inches(0.5)
    tf.margin_right = Inches(0.5)
    add_line(tf, "« Comment élaborer une architecture d'Intelligence Artificielle de type RAG "
             "souveraine, capable de transcender l'hétérogénéité des documents administratifs "
             "non structurés pour garantir une capitalisation fiable et une restitution "
             "sourcée de la mémoire institutionnelle ? »",
             20, NAVY, italic=True, font=SERIF, align=PP_ALIGN.CENTER, first=True)
    cap = add_textbox(s, Inches(1.0), Inches(5.2), Inches(11.33), Inches(1.0))
    add_line(cap, "Concilier innovation cognitive et gouvernance des données.",
             14, GREY, italic=True, font=SANS, align=PP_ALIGN.CENTER, first=True)
    notes(s, """Au croisement des impératifs opérationnels et des verrous techniques, la
    problématique : comment élaborer une architecture RAG souveraine capable de transcender
    l'hétérogénéité des documents non structurés pour garantir une capitalisation fiable et une
    restitution sourcée. Je lis la question lentement.""")


def slide_axes(prs):
    s = blank_slide(prs)
    header(s, "Problématique", "Axes de recherche et hypothèses")
    bullets(s, [
        "Axe 1 — Configuration des modèles open-source en local",
        ("Quelle combinaison embedding / génération / reranking pour un RAG juridique français ?", 1),
        ("H1 : HyDE et le reranking par cross-encoder améliorent les performances globales.", 1),
        "Axe 2 — Méthodologie d'évaluation objective",
        ("Comment mesurer précision de recherche et fidélité via RAGAS, face aux limites de "
         "l'évaluation automatique ?", 1),
        ("H2 : des modèles open-source locaux atteignent une qualité suffisante pour un usage "
         "juridique institutionnel.", 1),
    ], size=16)
    notes(s, """Deux axes. L'axe technique : quelle configuration de modèles en local, d'où H1
    sur l'apport de HyDE et du reranking. L'axe méthodologique : comment évaluer objectivement,
    d'où H2 sur la suffisance des modèles locaux. Je reviendrai sur la validation de H1 et H2
    dans les résultats.""")


# --------------------------------------------------------------------------- #
# Chapitre 1 — Revue de la littérature
# --------------------------------------------------------------------------- #
def slide_ch1_etat_art(prs):
    s = blank_slide(prs)
    header(s, "Chapitre 1 — Revue de la littérature", "L'état de l'art en cinq travaux")
    rows = [
        ["Travail", "Apport", "Brique du système"],
        ["Lewis et al. (2020)", "RAG : dissocier mémoire et raisonnement", "Architecture globale"],
        ["Reimers & Gurevych (2019)", "Bi-encodeur SBERT, similarité cosinus", "Embeddings (bge-m3)"],
        ["Gao et al. (2022)", "HyDE : combler le fossé lexical", "Recherche augmentée"],
        ["Nogueira & Cho (2019)", "Reranking cross-encoder", "ms-marco-MiniLM"],
        ["Es et al. (2023)", "RAGAS : évaluation par LLM-juge", "Protocole d'évaluation"],
    ]
    nrows, ncols = len(rows), len(rows[0])
    gt = s.shapes.add_table(nrows, ncols, Inches(0.7), Inches(1.9),
                            Inches(11.9), Inches(3.4)).table
    for i in range(nrows):
        for j in range(ncols):
            c = gt.cell(i, j)
            c.text = rows[i][j]
            p = c.text_frame.paragraphs[0]
            run = p.runs[0]
            run.font.size = Pt(13 if i == 0 else 12)
            run.font.name = SANS
            run.font.bold = (i == 0)
            run.font.color.rgb = WHITE if i == 0 else DARK
    add_line(add_textbox(s, Inches(0.7), Inches(5.6), Inches(11.9), Inches(0.8)),
             "Fil conducteur : chaque article résout la limite du précédent ; chaque brique de "
             "notre système en découle directement.", 14, GREY, italic=True, font=SANS, first=True)
    notes(s, """Je passe vite sur la théorie — le jury a lu le chapitre. L'essentiel : cinq
    travaux en fil conducteur. Lewis pose le RAG ; Reimers et Gurevych, les embeddings denses ;
    Gao, HyDE pour le fossé lexical ; Nogueira et Cho, le reranking ; Es, l'évaluation RAGAS.
    Chaque brique de mon système découle de l'un d'eux. Une phrase par ligne, je n'insiste pas.""")


def slide_ch1_positionnement(prs):
    s = blank_slide(prs)
    header(s, "Chapitre 1 · 1.5", "Positionnement de la recherche")
    bullets(s, [
        "Les briques existent dans la littérature, mais leur assemblage pour notre cas est inédit :",
        ("souveraineté — tout open-source et 100 % local, là où les articles utilisent des "
         "APIs cloud ;", 1),
        ("intégration combinée HyDE + reranking, rarement testée et mesurée ensemble ;", 1),
        ("spécialisation sur le juridique francophone marocain (Bulletin Officiel), peu couvert.", 1),
        "Apport : un pipeline souverain de bout en bout + un splitter juridique + une étude "
        "d'ablation et un protocole d'évaluation sur corpus réel.",
    ], size=16, top=Inches(1.9))
    notes(s, """Mon positionnement tient en trois originalités : la souveraineté totale en local,
    l'intégration combinée de HyDE et du reranking évaluée empiriquement, et la spécialisation
    juridique marocaine. L'apport concret : l'assemblage souverain, le splitter dédié et le
    protocole d'évaluation, que je détaille maintenant.""")


# --------------------------------------------------------------------------- #
# Chapitre 2 — Concepts, méthodes et outils
# --------------------------------------------------------------------------- #
def slide_ch2_archi(prs):
    s = blank_slide(prs)
    header(s, "Chapitre 2 · 2.1", "Architecture globale du pipeline")
    phases = [
        ("1. Ingestion", "Crawl + PDF", NAVY),
        ("2. Extraction", "Docling →\nMarkdown", NAVY),
        ("3. Indexation", "Splitter +\nbge-m3 · pgvector", ACCENT),
        ("4. Recherche", "HyDE + rerank\n+ génération", ACCENT),
        ("5. Conversation", "Historique\nglissant", GREY),
    ]
    n = len(phases)
    total_w = Inches(12.4)
    gap = Inches(0.18)
    bw = (total_w - gap * (n - 1)) / n
    left0 = Inches(0.45)
    top = Inches(2.7)
    for i, (titre, sub, color) in enumerate(phases):
        left = left0 + i * (bw + gap)
        b = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, bw, Inches(1.4))
        b.fill.solid()
        b.fill.fore_color.rgb = color
        b.line.fill.background()
        tf = b.text_frame
        tf.word_wrap = True
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        add_line(tf, titre, 13, WHITE, bold=True, font=SANS, align=PP_ALIGN.CENTER, first=True)
        add_line(tf, sub, 10, WHITE, font=SANS, align=PP_ALIGN.CENTER)
        if i < n - 1:
            arrow(s, left + bw - gap + Inches(0.005), top + Inches(0.48),
                  width=gap, height=Inches(0.45))
    bullets(s, [
        "Cinq phases séquentielles, découplées par une base de données centrale.",
        "Chaque composant implémente une interface abstraite → modèles interchangeables.",
        "Les phases 3 et 4 (en bleu) portent les contributions méthodologiques.",
    ], top=Inches(4.6), size=15)
    notes(s, """L'architecture comporte cinq phases séquentielles, communiquant par une base de
    données centrale, ce qui permet de réitérer une phase sans réimpacter les autres. Chaque
    composant implémente une interface abstraite, rendant les modèles interchangeables. Les
    phases 3 et 4 concentrent les contributions.""")


def slide_ch2_segmentation(prs):
    s = blank_slide(prs)
    header(s, "Chapitre 2 · 2.2", "Segmentation du corpus juridique marocain")
    bullets(s, [
        "Le Bulletin Officiel a une hiérarchie (Dahir > Loi > Titre > Chapitre > Article) "
        "encodée par la typographie, non balisée informatiquement.",
        "Défis : discontinuités de numérotation, densité variable des articles.",
        "Segmentation sémantique en deux étapes :",
        ("normalisation hiérarchique (_fix_heading_hierarchy + _KEYWORD_RULES) qui corrige les "
         "titres extraits par Docling ;", 1),
        ("découpage par heading puis sous-division récursive (chunk 1500, overlap 300).", 1),
        "Enrichissement par breadcrumbs : chaque chunk porte sa position hiérarchique complète "
        "→ vecteur enrichi et traçabilité de la réponse.",
    ], size=15)
    notes(s, """Section 2.2, ma contribution d'ingénierie. Le BO a une hiérarchie portée par la
    typographie, non balisée. Mon splitter la normalise d'abord — il corrige les titres mal
    extraits par Docling via un vocabulaire légal fixe — puis découpe par en-tête et subdivise
    récursivement à 1500 caractères avec 300 de recouvrement. Chaque chunk reçoit un fil
    d'Ariane, qui enrichit le vecteur et permet la restitution sourcée.""")


def slide_pratique_splitter(prs):
    s = blank_slide(prs)
    header(s, "Chapitre 2 · 2.2", "Algorithme de normalisation hiérarchique")
    code = (
        "def _classify(self, text):\n"
        "    s = text.strip()\n"
        "    if _ARTICLE_RE.match(s):\n"
        "        return None      # article → **gras**\n"
        "    for level, pat in _KEYWORD_RULES:\n"
        "        if pat.search(s):\n"
        "            return level # mot-clé légal → niveau\n"
        "    return -1            # titre libre → contexte"
    )
    code_box(s, Inches(0.7), Inches(1.7), Inches(6.2), Inches(2.9), code, size=12,
             title="moroccan_bo_splitter.py")
    # règles mots-clés → niveau (version lisible)
    rules = [
        ["Niveau", "Mots-clés (extrait)"],
        ["H1", "DAHIR · TEXTES GÉNÉRAUX · SOMMAIRE"],
        ["H2", "Dahir n° · Loi n° · Décret n° · Arrêté"],
        ["H3", "PARTIE · I. II. (romains)"],
        ["H4", "TITRE · 1. 2. (numéros)"],
        ["H5", "Chapitre · a. b. (lettres)"],
    ]
    gt = s.shapes.add_table(len(rules), 2, Inches(7.2), Inches(1.7),
                            Inches(5.4), Inches(2.4)).table
    for i, row in enumerate(rules):
        for j, val in enumerate(row):
            c = gt.cell(i, j)
            c.text = val
            run = c.text_frame.paragraphs[0].runs[0]
            run.font.size = Pt(11 if i else 12)
            run.font.name = SANS
            run.font.bold = (i == 0)
            run.font.color.rgb = WHITE if i == 0 else DARK
    # exemple avant / après
    ex = (
        "AVANT (Docling)        APRÈS (normalisé)\n"
        "# Loi n° 15-95     →   ## Loi n° 15-95   (instrument)\n"
        "# TITRE PREMIER    →   #### TITRE PREMIER (titre)\n"
        "## Article 78      →   **Article 78**     (gras)"
    )
    code_box(s, Inches(0.7), Inches(4.8), Inches(11.9), Inches(1.6), ex, size=11,
             title="Exemple de réécriture des niveaux")
    notes(s, """Concrètement, voici le cœur du splitter — la fonction _classify, que j'ai écrite.
    Elle décide, pour chaque titre : un article devient du gras ; un mot-clé légal reçoit un
    niveau fixe via la table de règles à droite ; un titre libre est inféré du contexte. En bas,
    un exemple : Docling sort des niveaux incohérents (Loi en H1, Article en titre), mon
    algorithme les réécrit selon la hiérarchie juridique et transforme l'article en gras. C'est
    ce qui rend les breadcrumbs fiables.""")


def slide_ch2_recherche(prs):
    s = blank_slide(prs)
    header(s, "Chapitre 2 · 2.3", "Stratégie de recherche : HyDE et reranking")
    bullets(s, [
        "Fossé lexical : la question en langage courant et le corpus normatif ont des "
        "vecteurs éloignés.",
        "HyDE : générer une réponse hypothétique employant le vocabulaire du corpus ; deux "
        "pools de récupération (requête + réponse hypothétique) fusionnés et dédoublonnés.",
        "Reranking : sur les ~20-30 candidats fusionnés, le cross-encoder "
        "ms-marco-MiniLM-L-6-v2 score chaque paire [requête, chunk].",
        "Retrieve-then-rerank : le bi-encodeur filtre largement, le cross-encoder affine "
        "précisément.",
        "HyDE et reranking sont les variables de l'hypothèse H1, isolées par ablation (3.2).",
    ], size=16)
    notes(s, """Section 2.3. Pour combler le fossé lexical, HyDE génère une réponse hypothétique
    et lance deux récupérations en parallèle, fusionnées. Puis le reranking : sur les candidats,
    le cross-encoder ms-marco score finement chaque paire. C'est le retrieve-then-rerank. Ces
    deux leviers sont précisément ce que l'étude d'ablation va isoler.""")


def slide_ch2_modeles(prs):
    s = blank_slide(prs)
    header(s, "Chapitre 2 · 2.4", "Modèles candidats et contraintes de sélection")
    bullets(s, [
        "Contrainte de souveraineté : modèles open-source, déployables localement (Ollama / "
        "HuggingFace), supportant le français.",
        ("Embedding de référence : bge-m3 (multilingue, 1024 dim, contexte 8192 tokens).", 1),
        ("Génération : qwen2.5:7b (référence), comparé à mistral:7b et gemma2:9b (chapitre 3).", 1),
        ("Reranking de référence : cross-encoder ms-marco-MiniLM-L-6-v2 (léger, exécutable CPU).", 1),
        ("Stockage : PostgreSQL + pgvector, index IVFFLAT (100 listes, 10 probes).", 1),
        "Critères juridiques : résistance au bruit, rejet explicite hors-corpus, fidélité aux "
        "sources.",
    ], size=15)
    notes(s, """Section 2.4. La souveraineté impose des modèles open-source locaux. L'embedding
    de référence est bge-m3, multilingue. La génération de référence est qwen2.5:7b, que je
    comparerai à mistral et gemma au chapitre 3. Le reranker est ms-marco, assez léger pour le
    CPU. Le stockage est pgvector. À noter : embedding et reranker sont fixés comme références ;
    seule la génération fait l'objet d'une comparaison empirique.""")


def slide_ch2_evaluation(prs):
    s = blank_slide(prs)
    header(s, "Chapitre 2 · 2.5", "Évaluation et construction du jeu de test")
    bullets(s, [
        "RAGAS mesure séparément retrieval et génération via quatre métriques :",
        ("faithfulness (ancrage, anti-hallucination) et answer relevancy (réponse à la "
         "question) → génération ;", 1),
        ("context precision (pertinence des chunks) et context recall (couverture) → retrieval.", 1),
        "Construction du jeu de test en trois étapes :",
        ("extraction de chunks seed (250-3000 car.) du corpus indexé ;", 1),
        ("génération synthétique des Q/R par gpt-4o-mini — distinct des modèles évalués "
         "(anti-biais d'auto-évaluation) ;", 1),
        ("validation humaine de chaque paire (formulation, exactitude, doublons).", 1),
    ], size=15)
    notes(s, """Section 2.5, le point que j'avais sous-développé. RAGAS fournit quatre métriques
    orthogonales : deux pour la génération — faithfulness et answer relevancy — deux pour le
    retrieval — context precision et recall. Le jeu de test est construit en trois étapes :
    extraction de chunks seed, génération synthétique par gpt-4o-mini — volontairement distinct
    des modèles évalués pour éviter le biais d'auto-évaluation — puis validation humaine de
    chaque paire. Protocole conforme à Es et al. (2023).""")


def slide_ch2_conversation(prs):
    s = blank_slide(prs)
    header(s, "Chapitre 2 · 2.6", "Gestion de la mémoire conversationnelle")
    bullets(s, [
        "Les questions de suivi (anaphores) imposent de maintenir un historique.",
        "Une accumulation illimitée dépasse la fenêtre de contexte du modèle (4096-8192 tokens).",
        "Stratégie : compression par résumé avec fenêtre glissante.",
        ("au-delà d'un seuil (~2000 tokens, estimés par len(message)//4), le LLM résume "
         "l'historique ancien ;", 1),
        ("les 4 derniers messages sont conservés intacts.", 1),
        "Aucune dépendance externe ; suffisant pour des sessions institutionnelles (10-20 échanges).",
    ], size=16)
    notes(s, """Section 2.6. Pour gérer les questions de suivi, je maintiens un historique. Mais
    il dépasserait vite la fenêtre du modèle. La stratégie : au-delà de 2000 tokens, le LLM
    résume l'historique ancien et conserve les 4 derniers messages intacts. Simple, sans
    dépendance externe, suffisant pour des sessions institutionnelles.""")


# --------------------------------------------------------------------------- #
# Chapitre 3 — Résultats
# --------------------------------------------------------------------------- #
def slide_ch3_protocole(prs):
    s = blank_slide(prs)
    header(s, "Chapitre 3 · 3.1", "Protocole expérimental")
    bullets(s, [
        "Jeu de test : 79 paires Q/R validées manuellement.",
        ("3 domaines (droit social, fiscal, sociétés) × 3 profils (factuel mono-article, "
         "synthèse multi-articles, hors-corpus).", 1),
        "Plan d'expérience : 12 configurations = 3 modèles de génération × HyDE (on/off) × "
        "reranking (on/off).",
        "Paramètres fixes : embedding bge-m3, reranker ms-marco, top_k = 6, seuil cosinus 0,5, "
        "pgvector IVFFLAT (100 listes, 10 probes).",
        "Juge RAGAS : gpt-4o-mini — externe au pipeline de production.",
        "Configuration de base : recherche dense par similarité cosinus, sans HyDE ni reranking.",
    ], size=15)
    notes(s, """Section 3.1. Le jeu de test : 79 paires stratifiées sur 3 domaines et 3 profils,
    dont un profil hors-corpus qui teste le refus. Le plan d'expérience couvre 12 configurations.
    Les paramètres non variables sont fixés : embedding bge-m3, reranker ms-marco, top_k 6, seuil
    0,5, pgvector. Le juge est gpt-4o-mini, hors production. La configuration de base est la
    recherche dense simple, sans HyDE ni reranking.""")


def slide_pratique_dataset(prs):
    s = blank_slide(prs)
    header(s, "Chapitre 3 · 3.1", "Exemples concrets du jeu de test")
    examples = [
        ("Factuel (mono-article)",
         "Q : « Quel montant le gouvernement peut-il dépenser en 2024 pour le fonds des titres "
         "identitaires électroniques ? »",
         "R : montant fixé par l'article 29.",
         "BO 7259-bis · art. 29 · 2023", NAVY),
        ("Synthèse (multi-articles)",
         "Q : « Comment fonctionne la fusion d'un OPCVM avec un autre OPCVM existant ? »",
         "R : fusion-absorption par la société de gestion (art. 78).",
         "BO 7462 · art. 78 · 2025", ACCENT),
        ("Hors-corpus (refus attendu)",
         "Q : « Comment puis-je obtenir un agrément pour 2024 ? »",
         "R : « L'extrait fourni ne contient pas d'information sur cette procédure. »",
         "BO 7410 · test anti-hallucination", GREY),
    ]
    top = Inches(1.75)
    h = Inches(1.65)
    gap = Inches(0.2)
    for i, (profil, q, r, src, color) in enumerate(examples):
        y = top + i * (h + gap)
        bar = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.7), y, Inches(0.12), h)
        bar.fill.solid()
        bar.fill.fore_color.rgb = color
        bar.line.fill.background()
        tf = add_textbox(s, Inches(0.95), y, Inches(11.6), h)
        add_line(tf, profil, 13, color, bold=True, font=SERIF, first=True)
        add_line(tf, q, 12, DARK, font=SANS, space_before=2)
        add_line(tf, r, 12, DARK, italic=True, font=SANS, space_before=1)
        add_line(tf, src, 10, GREY, font=SANS, space_before=1)
    notes(s, """Pour rendre cela concret, voici trois vraies questions du jeu de test, une par
    profil. Une question factuelle, dont la réponse tient dans un article — ici l'article 29 sur
    un montant budgétaire. Une question de synthèse, qui demande de combiner plusieurs articles —
    la fusion d'OPCVM, article 78. Et surtout une question hors-corpus : le système répond
    correctement qu'il n'a pas l'information, au lieu d'inventer. Ce profil teste explicitement
    l'anti-hallucination, exigence centrale en contexte juridique.""")


def slide_ch3_ablation(prs):
    s = blank_slide(prs)
    header(s, "Chapitre 3 · 3.2", "Étude d'ablation : impact du reranking et de HyDE")
    if (FIGURES / "fig_precision_config.png").exists():
        s.shapes.add_picture(str(FIGURES / "fig_precision_config.png"),
                             Inches(0.4), Inches(1.7), height=Inches(3.9))
    bullets(s, [
        "3.2.1 — Reranking seul :",
        ("+5 à 8 pts de précision du contexte, systématique sur les 3 modèles ;", 1),
        ("léger recul de la pertinence (contexte plus focalisé).", 1),
        "3.2.2 — HyDE seul :",
        ("dégrade la précision (~−8 pts) sur toutes les configs ;", 1),
        ("cause : écart de registre — le modèle paraphrase en français courant, pas dans le "
         "registre normatif du corpus.", 1),
    ], left=Inches(7.6), top=Inches(1.8), width=Inches(5.5), size=13)
    notes(s, """Section 3.2, l'étude d'ablation. Le reranking seul (3.2.1) apporte un gain
    systématique de précision, cinq à huit points, avec un léger recul de pertinence car le
    contexte est plus focalisé. HyDE seul (3.2.2), contre toute attente, dégrade la précision
    d'environ huit points : le modèle généraliste paraphrase en français courant au lieu de
    reproduire le registre normatif, ce qui éloigne le vecteur du corpus. C'est l'effet
    contre-intuitif central de mes résultats.""")


def slide_ch3_synthese(prs):
    s = blank_slide(prs)
    header(s, "Chapitre 3 · 3.2.3", "Synthèse de l'ablation et choix de configuration")
    if (FIGURES / "fig_scores_globaux.png").exists():
        s.shapes.add_picture(str(FIGURES / "fig_scores_globaux.png"),
                             Inches(0.4), Inches(1.7), height=Inches(3.5))
    runs = load_canonical_runs()
    rows = scores_table_12(runs)
    nrows, ncols = len(rows), len(rows[0])
    gt = s.shapes.add_table(nrows, ncols, Inches(0.9), Inches(5.4),
                            Inches(7.3), Inches(1.3)).table
    for i in range(nrows):
        for j in range(ncols):
            c = gt.cell(i, j)
            c.text = rows[i][j]
            run = c.text_frame.paragraphs[0].runs[0]
            run.font.size = Pt(10)
            run.font.name = SANS
            run.font.bold = (i == 0 or j == 0)
    cap = add_textbox(s, Inches(8.3), Inches(1.9), Inches(4.7), Inches(4.8))
    for i, line in enumerate([
        "hyde+rerank est marginalement le plus haut",
        "(qwen 0,794), mais HyDE est instable",
        "et coûte une 2ᵉ passe d'embedding (+30-50 % latence).",
        "",
        "Configuration retenue : sans HyDE + reranking",
        "— gain systématique, coût maîtrisé, meilleure précision",
        "(mistral 0,790 ; qwen 0,782).",
        "",
        "→ H1 validée partiellement.",
    ]):
        bold = line.startswith("Configuration") or line.startswith("→")
        add_line(cap, line, 13, NAVY if bold else DARK, bold=bold, font=SANS, first=(i == 0),
                 space_after=3)
    notes(s, """Section 3.2.3, la synthèse. Le tableau et la figure montrent les 12
    configurations. En score brut, HyDE plus reranking est marginalement le plus haut pour qwen,
    à 0,794. Mais HyDE est instable — il dégrade la précision — et coûte une seconde passe
    d'embedding, soit trente à cinquante pour cent de latence en plus. Je retiens donc la
    configuration sans HyDE avec reranking : un gain systématique, un coût maîtrisé, une
    meilleure précision. Si le jury demande pourquoi pas hyde+rerank, c'est exactement cette
    réponse. H1 est validée partiellement : le reranking confirme l'hypothèse, HyDE l'infirme.""")


def slide_ch3_generation(prs):
    s = blank_slide(prs)
    header(s, "Chapitre 3 · 3.3", "Comparaison des modèles de génération")
    if (FIGURES / "fig_models.png").exists():
        s.shapes.add_picture(str(FIGURES / "fig_models.png"),
                             Inches(0.5), Inches(1.7), height=Inches(3.4))
    bullets(s, [
        "Configuration fixée (sans HyDE + reranking) ; seule la génération varie.",
        "mistral:7b — meilleur score moyen (0,790) et pertinence (0,835) ;",
        "qwen2.5:7b — meilleure fidélité (0,709), précision élevée ;",
        "gemma2:9b — précision retrieval élevée mais fidélité/pertinence faibles (0,757).",
        "Précision du contexte stable (~0,85-0,88) : pipeline de retrieval commun.",
    ], top=Inches(5.3), size=14)
    notes(s, """Section 3.3. À configuration fixée, je compare les trois générateurs. Mistral a
    le meilleur score moyen et la meilleure pertinence. Qwen a la meilleure fidélité. Gemma a une
    bonne précision de retrieval mais une fidélité et une pertinence plus faibles. La précision
    du contexte est stable entre modèles, car le retrieval est commun.""")


def slide_ch3_interpretation(prs):
    s = blank_slide(prs)
    header(s, "Chapitre 3 · 3.3.3", "Interprétation : tension fidélité / pertinence")
    bullets(s, [
        "Aucun modèle ne maximise simultanément fidélité et pertinence — tension structurelle.",
        ("Profil conservateur — qwen2.5:7b : ancrage strict, fidélité 0,709, moins direct ;", 1),
        ("Profil assertif — mistral:7b : synthèse fluide, pertinence 0,835, ancrage plus fragile.", 1),
        "L'écart de fidélité (7 pts) n'est pas anodin en contexte juridique.",
        "Choix par profil de risque :",
        ("effet juridique direct (avis, conformité, instruction) → qwen2.5:7b ;", 1),
        ("exploration documentaire / formation → mistral:7b.", 1),
        "Modèles 7-9B locaux : scores 0,75-0,79 → H2 validée, sous arbitrage.",
    ], size=15)
    notes(s, """Section 3.3.3. L'interprétation clé : une tension structurelle entre fidélité et
    pertinence, qu'aucun modèle ne résout. Qwen est conservateur et fidèle ; mistral est assertif
    et pertinent. En droit, l'écart de fidélité compte. D'où un choix par profil de risque : qwen
    pour les usages à effet juridique direct, mistral pour l'exploration. Et globalement, des
    modèles locaux 7-9B atteignent 0,75 à 0,79 : H2 est validée, sous arbitrage.""")


# --------------------------------------------------------------------------- #
# Discussion & conclusion
# --------------------------------------------------------------------------- #
def slide_disc_hypotheses(prs):
    s = blank_slide(prs)
    header(s, "Discussion", "Validation des hypothèses et réponse à la problématique")
    bullets(s, [
        "H1 — validée partiellement : le reranking améliore les performances (précision +5-8 pts) ; "
        "HyDE seul les dégrade. Config optimale = sans HyDE + reranking.",
        "H2 — validée : des modèles open-source locaux atteignent un niveau exploitable, sous "
        "arbitrage fidélité / pertinence selon l'usage.",
        "Réponse à la problématique :",
        ("oui, une architecture RAG souveraine assure une capitalisation fiable et une "
         "restitution sourcée (breadcrumbs) de la mémoire institutionnelle ;", 1),
        ("la fidélité est maîtrisée par le choix de modèle ; l'exécution est 100 % locale.", 1),
    ], size=16)
    notes(s, """La discussion valide les hypothèses. H1 partiellement : le reranking confirme,
    HyDE infirme. H2 validée, sous arbitrage. Et surtout, je réponds à la problématique : oui,
    une architecture RAG souveraine permet une capitalisation fiable et une restitution sourcée,
    avec une fidélité maîtrisée par le choix de modèle et une exécution entièrement locale.""")


def slide_disc_fart(prs):
    s = blank_slide(prs)
    header(s, "Discussion", "Implications pour le déploiement au FART")
    bullets(s, [
        "Souveraineté : aucun appel cloud en production — PostgreSQL+pgvector, Ollama (GPU 8 Go), "
        "cross-encoder sur CPU ; seule l'évaluation RAGAS sollicite gpt-4o-mini.",
        "Choix de modèle de référence par profil de risque (qwen pour l'effet juridique direct).",
        "Transposition du BO (corpus de substitution) au corpus interne du FART :",
        ("composants transposables : bge-m3, ms-marco, génération, RAGAS ;", 1),
        ("deux écarts à combler : adapter le splitter à la terminologie FART, reconstruire le "
         "jeu de test interne ;", 1),
        ("effort estimé : une mission de l'ordre de 4 à 6 mois. [à confirmer]", 1),
    ], size=15)
    notes(s, """Les implications pour le FART. La souveraineté est garantie : aucun appel cloud
    en production, une infrastructure maîtrisable. Le modèle de référence se choisit par profil
    de risque. La transposition du BO au corpus interne est facilitée par l'équivalence
    structurelle : les composants principaux sont transposables, mais deux écarts demandent une
    adaptation — le splitter et le jeu de test. J'estime cet effort à quatre à six mois ; restez
    prudent sur ce chiffrage si on vous interroge.""")


def slide_conclusion(prs):
    s = blank_slide(prs)
    header(s, "Conclusion générale", "Contributions, limites et perspectives")
    bullets(s, [
        "Contributions : pipeline RAG souverain de bout en bout, splitter juridique, étude "
        "d'ablation HyDE/reranking, protocole d'évaluation RAGAS contrôlé.",
        "Limites : monolingue français ; 79 questions sur 3 domaines ; 3 modèles 7-9B ; juge "
        "externe gpt-4o-mini.",
        "Perspectives :",
        ("fine-tuning d'un modèle sur corpus juridique marocain (fidélité + registre HyDE) ;", 1),
        ("extension bilingue français-arabe, facilitée par bge-m3 ;", 1),
        ("déploiement effectif au FART ; réévaluation des modèles à 12-18 mois.", 1),
    ], size=15, height=Inches(4.3))
    big = add_textbox(s, Inches(0.7), Inches(6.2), Inches(11.9), Inches(0.6))
    add_line(big, "Merci de votre attention — je suis à votre disposition pour vos questions.",
             16, NAVY, bold=True, font=SERIF, align=PP_ALIGN.CENTER, first=True)
    notes(s, """En conclusion : les contributions sont le pipeline souverain complet, le splitter
    juridique, l'étude d'ablation et le protocole d'évaluation. Les limites, assumées :
    monolingue, jeu de test borné, modèles 7-9B, juge externe. Les perspectives : fine-tuning
    juridique marocain, extension bilingue facilitée par bge-m3, déploiement effectif au FART,
    et réévaluation périodique. Je vous remercie et reste à votre disposition.""")


def main():
    prs = Presentation()
    prs.slide_width = EMU_W
    prs.slide_height = EMU_H
    builders = [
        slide_titre, slide_plan, slide_intro, slide_probleme, slide_problematique, slide_axes,
        slide_ch1_etat_art, slide_ch1_positionnement,
        slide_ch2_archi, slide_ch2_segmentation, slide_pratique_splitter, slide_ch2_recherche,
        slide_ch2_modeles, slide_ch2_evaluation, slide_ch2_conversation,
        slide_ch3_protocole, slide_pratique_dataset, slide_ch3_ablation, slide_ch3_synthese,
        slide_ch3_generation, slide_ch3_interpretation,
        slide_disc_hypotheses, slide_disc_fart, slide_conclusion,
    ]
    for build in builders:
        build(prs)

    slides = list(prs.slides)
    total = len(slides)
    for i, s in enumerate(slides):
        if i == 0:
            continue
        footer(s, i + 1, total)

    prs.save(OUT)
    print(f"{total} slides → {OUT}")


if __name__ == "__main__":
    main()
